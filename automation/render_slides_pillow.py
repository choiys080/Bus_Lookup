import os
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

# Paths
pptx_path = r'd:\Antigravity\Bus_Lookup\ref\slides_temp.pptx'
output_dir = r'd:\Antigravity\Bus_Lookup\ref\slides'

if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Load presentation
try:
    prs = Presentation(pptx_path)
except Exception as e:
    print(f"Error loading PPTX: {e}")
    exit(1)

# Slide dimension (usually 960x540 or similar in pixels at 96dpi)
width, height = 1200, 675  # 16:9 ratio
bg_color = (255, 255, 255)
text_color = (15, 23, 42)  # Dark slate
accent_color = (0, 169, 122)  # B. Braun emerald

def render_slide(slide, slide_index):
    img = Image.new('RGB', (width, height), color=bg_color)
    draw = ImageDraw.Draw(img)

    # Try to load a font, fallback to default
    try:
        title_font = ImageFont.truetype("C:\\Windows\\Fonts\\malgun.ttf", 40)
        body_font = ImageFont.truetype("C:\\Windows\\Fonts\\malgun.ttf", 24)
    except:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()

    # Draw a small accent at the top
    draw.rectangle([0, 0, width, 5], fill=accent_color)

    y = 50
    margin = 50

    # Draw shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue

                # Simple layout: Title if at top, else body
                font = title_font if y < 150 else body_font

                # Draw text
                draw.text((margin, y), text, font=font, fill=text_color)

                # Move down (crude line spacing)
                y += 40 if font == title_font else 30

                if y > height - 50:
                    break

        # If it's a picture, we can't easily render it perfectly, but we could list it
        elif shape.shape_type == 13: # Picture
            draw.rectangle([width-150, height-150, width-50, height-50], outline=accent_color)
            draw.text((width-145, height-145), "[Image]", fill=accent_color)

    output_path = os.path.join(output_dir, f"slide_{slide_index}.png")
    img.save(output_path)
    print(f"Saved {output_path}")

# Process all slides
for i, slide in enumerate(prs.slides, 1):
    render_slide(slide, i)

print(f"Successfully rendered {len(prs.slides)} slides.")
