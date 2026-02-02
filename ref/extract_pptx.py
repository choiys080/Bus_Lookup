from pptx import Presentation

prs = Presentation(r'd:\Antigravity\Bus_Lookup\ref\비브라운코리아_액티비티 앱0130.pptx')

with open(r'd:\Antigravity\Bus_Lookup\ref\pptx_content.txt', 'w', encoding='utf-8') as f:
    for slide_num, slide in enumerate(prs.slides, 1):
        f.write('\n' + '='*50 + '\n')
        f.write(f'SLIDE {slide_num}\n')
        f.write('='*50 + '\n')

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ''.join([run.text for run in para.runs])
                    if text.strip():
                        f.write(text.strip() + '\n')

print('Saved to pptx_content.txt')
